#imports for the functions defined here
from ast import Try
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import *
from pptx.enum.dml import MSO_THEME_COLOR
import pptx
from PIL import ImageColor
from tkinter.messagebox import showerror
from Backend.Custom_Errors import ColorSelectionError

#Creates a textbox in presentation
def textBox(top: float,left: float,width: float,height: float,shapes: pptx.shapes.shapetree.SlideShapes,text: str,
is_bold=False, is_italics=False,is_underline=False,font_size=12,font_name="Arial",font_color="#000000",fill_color=None,
outline_color=None,outline_width=None,text_align="Left",task_level=1):
    
    textbox = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),Inches(width),Inches(height))
    if fill_color==None:
        textbox.fill.background()
    else:
        try:
            fill_color_r=ImageColor.getcolor(fill_color,"RGB")[0]
            fill_color_g=ImageColor.getcolor(fill_color,"RGB")[1]
            fill_color_b=ImageColor.getcolor(fill_color,"RGB")[2]
        except:
            showerror("Text Fill Color Not Selected",
            f"Please ensure text fill color for task {task_level} is selected.")
            raise ColorSelectionError("Text Fill Color Not Selected")
        try:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(fill_color_r,fill_color_g,fill_color_b)
        except:
            print("Enter a valid hex color code. Ex. #FFFF00")

    if outline_color==None:
        textbox.line.fill.background()
    else:
        try:
            outline_color_r=ImageColor.getcolor(outline_color,"RGB")[0]
            outline_color_g=ImageColor.getcolor(outline_color,"RGB")[1]
            outline_color_b=ImageColor.getcolor(outline_color,"RGB")[2]
        except:
            showerror("Text Outline Color Not Selected",
            f"Please ensure text outline color for task {task_level} is selected.")
            raise ColorSelectionError("Text Outline Color Not Selected")
        try:
            textbox.line.color.rgb = RGBColor(outline_color_r,outline_color_g,outline_color_b)
        except:
            print("Enter a valid hex color code. Ex. #FFFF00")
    
    if outline_width!=None:
        textbox.line.width = Pt(outline_width)
    
    p=textbox.text_frame.paragraphs[0]
    
    if text_align.lower()=="right":
        p.alignment = PP_ALIGN.RIGHT
    elif text_align.lower()=="center":
        p.alignment = PP_ALIGN.CENTER
    else:
        p.alignment = PP_ALIGN.LEFT
    
    run = p.add_run()
    run.text= text
    run.font.bold=is_bold
    run.font.size = Pt(font_size)
    run.font.italic = is_italics
    run.font.underline = is_underline
    run.font.name = font_name
    try:
        font_color_r=ImageColor.getcolor(font_color,"RGB")[0]
        font_color_g=ImageColor.getcolor(font_color,"RGB")[1]
        font_color_b=ImageColor.getcolor(font_color,"RGB")[2]
    except:
        showerror("Text Font Color Not Selected",
        f"Please ensure text font color for task {task_level} is selected.")
        raise ColorSelectionError("Text Font Color Not Selected")
    try: 
         run.font.color.rgb = RGBColor(font_color_r,font_color_g,font_color_b)#.lstrip('#'))
    except:
        print("Enter a valid hex color code. Ex. #FFFF00")

#creates a shape for gantt in the presentation
def ganttShape(top: float,left: float, width: float, height: float,shapes: pptx.shapes.shapetree.SlideShapes,
fill_color=None,outline_color=None,outline_width=None,shape_name="Pentagon",brightness=0,task_level=1):
    
    if shape_name.lower()=="rectangle":
        ganttshape = shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left), Inches(top), Inches(width), Inches(height))
    elif shape_name.lower()=="chevron":
        ganttshape = shapes.add_shape(MSO_SHAPE.CHEVRON,Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        ganttshape = shapes.add_shape(MSO_SHAPE.PENTAGON,Inches(left), Inches(top), Inches(width), Inches(height))
    
    ganttshape.fill.solid()
    
    if fill_color==None:
        ganttshape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        ganttshape.fill.fore_color.brightness = brightness
    else:
        try:
            fill_color_r=ImageColor.getcolor(fill_color,"RGB")[0]
            fill_color_g=ImageColor.getcolor(fill_color,"RGB")[1]
            fill_color_b=ImageColor.getcolor(fill_color,"RGB")[2]
        except:
            showerror("Shape Fill Color Not Selected",
            f"Please ensure shape fill color for task {task_level} is selected.")
            raise ColorSelectionError("Shape Fill Color Not Selected")
        try:
            ganttshape.fill.fore_color.rgb = RGBColor(fill_color_r,fill_color_g,fill_color_b)
        except:
            print("Enter a valid hex color code. Ex. #FFFF00")
    
    
    if outline_color==None:
        ganttshape.line.fill.background()
    else:
        try:
            outline_color_r=ImageColor.getcolor(outline_color,"RGB")[0]
            outline_color_g=ImageColor.getcolor(outline_color,"RGB")[1]
            outline_color_b=ImageColor.getcolor(outline_color,"RGB")[2]
        except:
            showerror("Shape Outline Color Not Selected",
            f"Please ensure shape outline color for task {task_level} is selected.")
            raise ColorSelectionError("Shape Outline Color Not Selected")
        try:
            ganttshape.line.color.rgb = RGBColor(outline_color_r,outline_color_g,outline_color_b)
        except:
            print("Enter a valid hex color code. Ex. #FFFF00")
    
    if outline_width!=None:
        ganttshape.line.width = Pt(outline_width)

def milestoneShape(left: float, top: float, width: float, height: float, shapes: pptx.shapes.shapetree.SlideShapes,
fill_color='#000000',outline_color=None,outline_width=None,shape_name="Diamond",brightness=0,task_level=1):
    if shape_name.lower()=='star':
        milestoneshape = shapes.add_shape(MSO_SHAPE.STAR_5_POINT,Inches(left),Inches(top), Inches(width), Inches(height))
    elif shape_name.lower()=="square":
        milestoneshape = shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left),Inches(top), Inches(width), Inches(height))
    elif shape_name.lower()=='triangle':
        milestoneshape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left),Inches(top), Inches(width), Inches(height))
    elif shape_name.lower()=="circle":
        milestoneshape = shapes.add_shape(MSO_SHAPE.OVAL,Inches(left),Inches(top), Inches(width), Inches(height))
    else:
        milestoneshape = shapes.add_shape(MSO_SHAPE.DIAMOND,Inches(left),Inches(top), Inches(width), Inches(height))
    milestoneshape.fill.solid()
    
    if fill_color==None:
        milestoneshape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        milestoneshape.fill.fore_color.brightness = brightness
    else:
        try:
            fill_color_r=ImageColor.getcolor(fill_color,"RGB")[0]
            fill_color_g=ImageColor.getcolor(fill_color,"RGB")[1]
            fill_color_b=ImageColor.getcolor(fill_color,"RGB")[2]
        except:
            showerror("Shape Fill Color Not Selected",
            f"Please ensure shape fill color for task {task_level} is selected.")
            raise ColorSelectionError("Shape Fill Color Not Selected")
        try:
            milestoneshape.fill.fore_color.rgb = RGBColor(fill_color_r,fill_color_g,fill_color_b)
        except:
            print("Enter a valid hex color code. Ex. #FFFF00")
    
    if outline_color==None:
        milestoneshape.line.fill.background()
    else:
        try:
            outline_color_r=ImageColor.getcolor(outline_color,"RGB")[0]
            outline_color_g=ImageColor.getcolor(outline_color,"RGB")[1]
            outline_color_b=ImageColor.getcolor(outline_color,"RGB")[2]
        except:
            showerror("Shape Outline Color Not Selected",
            f"Please ensure shape outline color for task {task_level} is selected.")
            raise ColorSelectionError("Shape Outline Color Not Selected")
        
        try:
            milestoneshape.line.color.rgb = RGBColor(outline_color_r,outline_color_g,outline_color_b)
        except:
            print("Enter a valid hex color code. Ex. #FFFF00")
    
    if outline_width!=None:
        milestoneshape.line.width = Pt(outline_width)    