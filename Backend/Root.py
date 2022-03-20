#imports for the application file
# Work to improve the code with additional tables and functionality
'''
Things in Mind:
1. Adding Slide Title
2. Adding Legend to top right
3. Adding Task Timeline Number If Possible
4. Adding Condition to table type based on selection
5. Changing Condition to Create a New Slide if there is change in the Slide Title
6. Change show previous slide function
'''

from Backend.Gantt_Shapes import *
from Backend.Gantt_Tables import *
from Backend.Util_Functions import *
from Backend.Custom_Errors import *
from Frontend.strings import months_dict
from tkinter.messagebox import showerror
import pandas as pd
import numpy as np
import datetime as dt
import pptx

# Creates a dataframe with required columns for creating gantt charts
def data_setup_for_gantt(data: pd.DataFrame,task_name: str,task_duration: str,
task_level: str,gantt_start_date: str,task_start_date: str,gantt_duration: str, slide_title: str):
    # Creates a new data frame
    gantt_data = data

    # Renames the column to prepare further operations
    gantt_data.rename(columns={
        task_name:"task_name",
        task_duration:"task_duration",
        task_level:"task_level",
        gantt_start_date:"gantt_start_date",
        task_start_date:"task_start_date",
        gantt_duration:"gantt_duration",
        slide_title:"slide_title"
    },inplace=True)

    # Changes the type of the data incase the type is not good
    try:
        gantt_data['task_duration']=gantt_data['task_duration'].astype(float)
    except:
        showerror("Invalid Task Duration",
        "Check the selection of Task Duration. Please make sure task duration is a numeric value and is in months.")
    try:
        gantt_data['gantt_duration']=gantt_data['gantt_duration'].astype(float)
    except:
        showerror("Invalid Gantt Duration",
        "Check the selection of Gantt Duration. Please make sure gantt duration is a numeric value and is in months.")
    try:
        gantt_data['task_level']=gantt_data['task_level'].astype(int)
    except:
        showerror("Invalid Task Level",
        "Check the selection of Task Level. Please make sure that task level is a numeric value between 1 to 5")
    try:
        gantt_data['task_name']=gantt_data['task_name'].astype("string")
    except:
        showerror("Invalid Task Name",
        "Check the selection of Task Name. Please make sure that the task name is text")
    try:
        gantt_data['slide_title']=gantt_data['slide_title'].astype("string")
    except:
        showerror("Invalid Slide Title",
        "Check the selection of Slide Title. Please make sure that the slide title is text")
    try:
        gantt_data['gantt_start_date']=pd.to_datetime(gantt_data['gantt_start_date'])
    except:
        showerror("Invalid Gantt Start Date",
        "Check the selection for Gantt Start Date. Please make sure that the 'Gantt Start Date' is in date format")
    try:
        gantt_data['task_start_date']=pd.to_datetime(gantt_data['task_start_date'])
    except:
        showerror("Invalide Task Start Date",
        "Check the seelction for Task Start Date. Please make sure that the 'Task Start Date' is in date format")
    # Create helper columns and other mathematical calculations
    gantt_data['Start Calibration']=gantt_data['task_start_date']-gantt_data['gantt_start_date']
    gantt_data['Start Calibration']=gantt_data['Start Calibration']/np.timedelta64(1,"M")
    gantt_data['Start Calibration']=gantt_data['Start Calibration'].astype(int)
    return gantt_data

# Creates the gantt slides and presentation
def create_gantt_slides(prs: pptx.Presentation,gantt_data: pd.DataFrame, timeline_input: dict,
align_tl: list,shape_type_tl: list,shape_color_tl:list,font_prop_tl: list,font_style_tl:list,
font_color_tl: list,font_size_tl: list, shape_height_tl: list, shape_height_metric_tl: list):
    # Decoding the inputs from the dict

    slide_layout_name = timeline_input.get('Layout')[0]
    create_timeline = timeline_input.get('Timeline')[0]
    create_milestones = timeline_input.get('Milestone')[0]
    table_granularity = timeline_input.get('Granularity')[0]

    table_left = timeline_input.get('Left')[0]
    table_left_metric = timeline_input.get('Left')[1]
    table_top = timeline_input.get("Top")[0]
    table_top_metric = timeline_input.get("Top")[1]
    table_width = timeline_input.get("Width")[0]
    table_width_metric = timeline_input.get("Width")[1]
    table_height  = timeline_input.get("Height")[0]
    table_height_metric = timeline_input.get("Height")[1]

    gantt_duration = int(gantt_data['gantt_duration'][0])
    gantt_start_date = gantt_data['gantt_start_date'][0]
    # Changing the data type for table properties
    try:
        table_left = getInches(float(table_left),table_left_metric)
        table_top = getInches(float(table_top),table_top_metric)
        table_width = getInches(float(table_width),table_width_metric)
        table_height = getInches(float(table_height),table_height_metric)
    except:
        showerror("Enter valid timeline properties",
        f"Timeline properties must be numeric. Cross check these inputs \n{table_left}\n{table_top}\n{table_width}\n{table_height}")
    
    try:
        for i in range(len((shape_height_tl))):
            shape_height_tl[i] = getInches(float(shape_height_tl[i]),shape_height_metric_tl[i])
    except:
        showerror("Enter valid task shape height input",
        f"Please check the shape height input and ensure it is numeric.")
    # Defines the slide layout
    slide_layout = prs.slide_layouts.get_by_name(slide_layout_name,default="Blank")
    # All the variables required for Gantt Automation Operation
    width_per_col = table_width/gantt_duration
    max_slide_height = 7.5
    table_columns = gantt_duration
    table_rows = 2
    height_measure = table_top
    slide_counter = 0
    task_height = 0.2
    
    slide_title_list = list(gantt_data['slide_title'].unique())
    slide_title_index = []
    for i in range(len(slide_title_list)):
        slide_title_index.append(i)
    slide_title_map = dict(zip(slide_title_list,slide_title_index))
    gantt_data['slide_title_number'] = gantt_data['slide_title'].map(slide_title_map)
    # Creates the first slide with table
    slide = [prs.slides.add_slide(slide_layout)]
    create_slide_title(slide[slide_counter].shapes,timeline_input,gantt_data['slide_title'][0])
    slide_title_number = gantt_data['slide_title_number'][0]
    create_legend(slide[slide_counter].shapes,shape_type_tl,shape_color_tl,table_top)
    if create_timeline:
        table = [createTable(table_rows,int(table_columns),table_left,table_top,table_width,table_height,slide[0])]
        create_timeline_function(table[slide_counter],table_granularity,gantt_start_date,gantt_duration,timeline_input)
        if create_milestones:
            additional_height = create_timeline_milestones(slide[slide_counter].shapes,timeline_input,
                    table_top+table_height*table_rows-0.2,
                    table_left,
                    gantt_start_date,
                    width_per_col)
            task_top = table_top+table_height*table_rows + additional_height
    else:
        task_top = table_top+table_height*table_rows-0.2
    # Create the main loop to do all the operation
    for i in range(len(gantt_data)):
        height_measure=task_top+task_height
        if (height_measure>=(max_slide_height-0.5)) or (gantt_data['slide_title_number'][i]>slide_title_number):
            slide_counter=slide_counter+1
            new_slide = prs.slides.add_slide(slide_layout)
            slide.append(new_slide)
            height_measure=table_top
            create_slide_title(slide[slide_counter].shapes,timeline_input,gantt_data['slide_title'][i])
            slide_title_number = gantt_data['slide_title_number'][i]
            create_legend(slide[slide_counter].shapes,shape_type_tl,shape_color_tl,table_top)
            if create_timeline:
                new_table = createTable(table_rows,int(table_columns),table_left,table_top,table_width,table_height,slide[slide_counter])
                table.append(new_table)
                create_timeline_function(table[slide_counter],table_granularity,gantt_start_date,gantt_duration,timeline_input)
                if create_milestones:
                    additional_height = create_timeline_milestones(slide[slide_counter].shapes,timeline_input,
                    table_top+table_height*table_rows-0.2,
                    table_left,
                    gantt_start_date,
                    width_per_col)
                    task_top = table_top+table_height*table_rows + additional_height
            else:
                task_top = table_top+table_height*table_rows-0.2
        if gantt_data['task_level'][i]==1:
            try:
                shape_left=table_left+gantt_data['Start Calibration'][i]*width_per_col
                shape_width = gantt_data['task_duration'][i]*width_per_col
                ganttShape(
                    top=task_top,
                    left=shape_left,
                    width=shape_width,
                    height=shape_height_tl[0],
                    shapes=slide[slide_counter].shapes,
                    shape_name=shape_type_tl[0],
                    fill_color=shape_color_tl[0]
                )
                text_left = set_textbox_shape_alignment(align_tl[0],shape_left=shape_left,shape_width=shape_width)
                prop=set_font_properties(font_prop_tl[0])

                if align_tl[0].lower()=="bottom of shape":
                    task_top=task_top+task_height+0.05
                
                textBox(
                    top=task_top,
                    left=text_left,
                    width=table_left+1,
                    height=task_height,
                    shapes=slide[slide_counter].shapes,
                    is_bold=prop[0],
                    is_italics=prop[1],
                    is_underline=prop[2],
                    font_name=font_style_tl[0],
                    font_color=font_color_tl[0],
                    font_size=int(font_size_tl[0]),
                    text=str(gantt_data['task_name'][i]),
                    text_align="Left"
                )
                task_top=task_top+task_height+0.1
            except ColorSelectionError:
                break
        elif gantt_data['task_level'][i]==2:
            try:
                shape_left=table_left+gantt_data['Start Calibration'][i]*width_per_col
                shape_width = gantt_data['task_duration'][i]*width_per_col
                ganttShape(
                    top=task_top,
                    left=shape_left,
                    width=shape_width,
                    height=shape_height_tl[1],
                    shapes=slide[slide_counter].shapes,
                    shape_name=shape_type_tl[1],
                    fill_color=shape_color_tl[1]
                )
                text_left = set_textbox_shape_alignment(align_tl[1],shape_left=shape_left,shape_width=shape_width)
                prop=set_font_properties(font_prop_tl[1])

                if align_tl[1].lower()=="bottom of shape":
                    task_top=task_top+task_height+0.05
                
                textBox(
                    top=task_top,
                    left=text_left,
                    width=table_left+1,
                    height=task_height,
                    shapes=slide[slide_counter].shapes,
                    is_bold=prop[0],
                    is_italics=prop[1],
                    is_underline=prop[2],
                    font_name=font_style_tl[1],
                    font_color=font_color_tl[1],
                    font_size=int(font_size_tl[1]),
                    text=str(gantt_data['task_name'][i]),
                    text_align="Left"
                )
                task_top=task_top+task_height+0.1
            except ColorSelectionError:
                break
        elif gantt_data['task_level'][i]==3:
            try:
                shape_left=table_left+gantt_data['Start Calibration'][i]*width_per_col
                shape_width = gantt_data['task_duration'][i]*width_per_col
                ganttShape(
                    top=task_top,
                    left=shape_left,
                    width=shape_width,
                    height=shape_height_tl[2],
                    shapes=slide[slide_counter].shapes,
                    shape_name=shape_type_tl[2],
                    fill_color=shape_color_tl[2]
                )
                text_left = set_textbox_shape_alignment(align_tl[2],shape_left=shape_left,shape_width=shape_width)
                prop=set_font_properties(font_prop_tl[2])

                if align_tl[2].lower()=="bottom of shape":
                    task_top=task_top+task_height+0.05
                
                textBox(
                    top=task_top,
                    left=text_left,
                    width=table_left+1,
                    height=task_height,
                    shapes=slide[slide_counter].shapes,
                    is_bold=prop[0],
                    is_italics=prop[1],
                    is_underline=prop[2],
                    font_name=font_style_tl[2],
                    font_color=font_color_tl[2],
                    font_size=int(font_size_tl[2]),
                    text=str(gantt_data['task_name'][i]),
                    text_align="Left"
                )
                task_top=task_top+task_height+0.1
            except ColorSelectionError:
                break
        elif gantt_data['task_level'][i]==4:
            try:
                shape_left=table_left+gantt_data['Start Calibration'][i]*width_per_col
                shape_width = gantt_data['task_duration'][i]*width_per_col
                milestoneShape(
                    top=task_top,
                    left=shape_left,
                    width=shape_height_tl[3],
                    height=shape_height_tl[3],
                    shapes=slide[slide_counter].shapes,
                    shape_name=shape_type_tl[3],
                    fill_color=shape_color_tl[3]
                )
                text_left = set_textbox_shape_alignment(align_tl[3],shape_left=shape_left,shape_width=shape_width)
                prop=set_font_properties(font_prop_tl[3])

                if align_tl[3].lower()=="bottom of shape":
                    task_top=task_top+task_height+0.05
                
                textBox(
                    top=task_top,
                    left=text_left,
                    width=table_left+1,
                    height=task_height,
                    shapes=slide[slide_counter].shapes,
                    is_bold=prop[0],
                    is_italics=prop[1],
                    is_underline=prop[2],
                    font_name=font_style_tl[3],
                    font_color=font_color_tl[3],
                    font_size=int(font_size_tl[3]),
                    text=str(gantt_data['task_name'][i]),
                    text_align="Left"
                )
                task_top=task_top+task_height+0.1
            except ColorSelectionError:
                break
        else:
            try:
                shape_left=table_left+gantt_data['Start Calibration'][i]*width_per_col
                shape_width = gantt_data['task_duration'][i]*width_per_col
                milestoneShape(
                    top=task_top,
                    left=shape_left,
                    width=shape_height_tl[4],
                    height=shape_height_tl[4],
                    shapes=slide[slide_counter].shapes,
                    shape_name=shape_type_tl[4],
                    fill_color=shape_color_tl[4]
                )
                text_left = set_textbox_shape_alignment(align_tl[4],shape_left=shape_left,shape_width=shape_width)
                prop=set_font_properties(font_prop_tl[4])

                if align_tl[4].lower()=="bottom of shape":
                    task_top=task_top+task_height+0.05
                
                textBox(
                    top=task_top,
                    left=text_left,
                    width=table_left+1,
                    height=task_height,
                    shapes=slide[slide_counter].shapes,
                    is_bold=prop[0],
                    is_italics=prop[1],
                    is_underline=prop[2],
                    font_name=font_style_tl[4],
                    font_color=font_color_tl[4],
                    font_size=int(font_size_tl[4]),
                    text=str(gantt_data['task_name'][i]),
                    text_align="Left"
                )
                task_top=task_top+task_height+0.1
            except ColorSelectionError:
                break   
    return prs

# Function to create timeline milestones
def create_timeline_milestones(shapes,timeline_input: dict, top, table_left, gantt_start_date, widthpercol):
    milestone_name = "Milestone "
    new_top = top
    for i in range(1,5):
        if timeline_input.get(milestone_name+str(i))[0]:
            year = int(timeline_input.get(milestone_name+str(i))[7])
            month = int(months_dict.get(timeline_input.get(milestone_name+str(i))[6]))
            shape_left = table_left + ((dt.datetime(year,month,1)-gantt_start_date)/np.timedelta64(1,"M"))*widthpercol
            milestoneShape(
                left = shape_left,
                top = top-0.1,
                width = 0.2,
                height = 0.2,
                shapes = shapes,
                outline_width =1,
                shape_name = timeline_input.get(milestone_name+str(i))[1],
                fill_color = timeline_input.get(milestone_name+str(i))[2],
                outline_color = timeline_input.get(milestone_name+str(i))[3]
            )
            if timeline_input.get(milestone_name+str(i))[4].lower()=="left of shape":
                textbox_left = shape_left - table_left
            elif timeline_input.get(milestone_name+str(i))[4].lower()=="right of shape":
                textbox_left = shape_left + 0.3
            else:
                textbox_left = shape_left
                new_top = top + 0.2
            textBox(
                left = textbox_left,
                top = new_top-0.1,
                width = 1.5,
                height = 0.2,
                shapes = shapes,
                text = timeline_input.get(milestone_name+str(i))[5],
                font_name = timeline_input.get("Milestone Properties")[0][0],
                font_size = timeline_input.get("Milestone Properties")[0][1],
                is_bold = timeline_input.get("Milestone Properties")[0][2],
                is_italics = timeline_input.get("Milestone Properties")[0][3],
                font_color = timeline_input.get("Milestone Properties")[1],
                fill_color = timeline_input.get("Milestone Properties")[2],
                text_align = timeline_input.get("Milestone Properties")[3]
            )
    return new_top-top


# Function to create slide title
def create_slide_title(shapes,timeline_input: dict, slide_title: str):
    textBox(
        top = 0.5,
        left = 0.8,
        width = 10,
        height = 0.5,
        text = slide_title,
        shapes = shapes,
        font_name = timeline_input.get('Title')[0][0],
        font_size = int(timeline_input.get('Title')[0][1]),
        is_bold = timeline_input.get('Title')[0][2],
        is_italics = timeline_input.get('Title')[0][3],
        font_color = timeline_input.get('Title')[1],
        fill_color = timeline_input.get('Title')[2],
        text_align = timeline_input.get('Title')[3]
    )

# Function to create legend
def create_legend(shapes,shape_type_tl,shape_color_tl,table_top):
    pass

# Function to create timeline table
def create_timeline_function(table, table_granularity, gantt_start_date, gantt_duration, timeline_input: dict):
    if table_granularity=="Month":
        simpleMonthlyTable(
            table = table,
            start_month = gantt_start_date,
            no_of_months = int(gantt_duration),
            font_name_year = timeline_input.get('Row 1 Properties')[0][0],
            font_size_year = timeline_input.get('Row 1 Properties')[0][1],
            is_bold_year = timeline_input.get('Row 1 Properties')[0][2],
            is_italics_year = timeline_input.get('Row 1 Properties')[0][3],
            font_color_year = timeline_input.get('Row 1 Properties')[1],
            fill_color_year = timeline_input.get('Row 1 Properties')[2],
            text_align_year = timeline_input.get('Row 1 Properties')[3],
            font_name_month = timeline_input.get('Row 2 Properties')[0][0],
            font_size_month = timeline_input.get('Row 2 Properties')[0][1],
            is_bold_month = timeline_input.get('Row 2 Properties')[0][2],
            is_italics_month = timeline_input.get('Row 2 Properties')[0][3],
            font_color_month = timeline_input.get('Row 2 Properties')[1],
            fill_color_month = timeline_input.get('Row 2 Properties')[2],
            text_align_month = timeline_input.get('Row 2 Properties')[3],
            )
    elif table_granularity=="Quarter":
        simpleQuarterlyTable(
            table = table,
            start_month = gantt_start_date,
            no_of_months = int(gantt_duration),
            font_name_year = timeline_input.get('Row 1 Properties')[0][0],
            font_size_year = timeline_input.get('Row 1 Properties')[0][1],
            is_bold_year = timeline_input.get('Row 1 Properties')[0][2],
            is_italics_year = timeline_input.get('Row 1 Properties')[0][3],
            font_color_year = timeline_input.get('Row 1 Properties')[1],
            fill_color_year = timeline_input.get('Row 1 Properties')[2],
            text_align_year = timeline_input.get('Row 1 Properties')[3],
            font_name_month = timeline_input.get('Row 2 Properties')[0][0],
            font_size_month = timeline_input.get('Row 2 Properties')[0][1],
            is_bold_month = timeline_input.get('Row 2 Properties')[0][2],
            is_italics_month = timeline_input.get('Row 2 Properties')[0][3],
            font_color_month = timeline_input.get('Row 2 Properties')[1],
            fill_color_month = timeline_input.get('Row 2 Properties')[2],
            text_align_month = timeline_input.get('Row 2 Properties')[3],
            )
    else:
        simpleSemiAnnualTable(
            table = table,
            start_month = gantt_start_date,
            no_of_months = int(gantt_duration),
            font_name_year = timeline_input.get('Row 1 Properties')[0][0],
            font_size_year = timeline_input.get('Row 1 Properties')[0][1],
            is_bold_year = timeline_input.get('Row 1 Properties')[0][2],
            is_italics_year = timeline_input.get('Row 1 Properties')[0][3],
            font_color_year = timeline_input.get('Row 1 Properties')[1],
            fill_color_year = timeline_input.get('Row 1 Properties')[2],
            text_align_year = timeline_input.get('Row 1 Properties')[3],
            font_name_month = timeline_input.get('Row 2 Properties')[0][0],
            font_size_month = timeline_input.get('Row 2 Properties')[0][1],
            is_bold_month = timeline_input.get('Row 2 Properties')[0][2],
            is_italics_month = timeline_input.get('Row 2 Properties')[0][3],
            font_color_month = timeline_input.get('Row 2 Properties')[1],
            fill_color_month = timeline_input.get('Row 2 Properties')[2],
            text_align_month = timeline_input.get('Row 2 Properties')[3],
            )

# Function to align textboxes based on the selection
def set_textbox_shape_alignment(align: str, shape_left: float, shape_width: float):
    max_width = 13.3
    if align.lower()=="left of slide":
        return 0.3
    elif align.lower()=="right of slide":
        return 12.0
    elif align.lower()=="left of shape":
        return shape_left-2.0
    elif align.lower()=="right of shape":
        if max_width - shape_left - shape_width>2.0:
            return shape_width+shape_left+0.2
        return shape_width+shape_left-2.0
    else:
        return shape_left

# Function to set font properties in text box based on the selection
def set_font_properties(prop: str)->list():
    if prop.lower()=="regular":
        return [False,False,False]
    elif prop.lower()=="bold":
        return [True,False,False]
    elif prop.lower()=="italic":
        return [False,True,False]
    elif prop.lower()=="underline":
        return [False,False,True]
    elif prop.lower()=="bold italic":
        return [True,True,False]
    elif prop.lower()=="bold underline":
        return [True,False,True]
    elif prop.lower()=="underline italic":
        return [False,True,True]
    else:
        return [False,False,False]
    