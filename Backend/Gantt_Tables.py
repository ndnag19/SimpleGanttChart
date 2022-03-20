#imports for the functions defined here
import pandas as pd
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import *
from dateutil.relativedelta import relativedelta
import pptx
from tkinter.messagebox import showerror
from PIL import ImageColor

#Create a Table
def createTable(rows: int,cols: int, left: float,top: float,width: float,height: float, slide):
    shape=slide.shapes.add_table(rows,cols,Inches(left), Inches(top), Inches(width), Inches(height))
    return shape.table


#Create a Monthly Table
def simpleMonthlyTable(table: pptx.table.Table,start_month: pd._libs.tslibs.timestamps.Timestamp, 
no_of_months: int, font_size_year=12, font_size_month=11, font_color_year="#FFFFFF", font_color_month="#FFFFFF", 
font_name_year="Arial", font_name_month = "Arial", fill_color_year=None, fill_color_month=None,
text_align_year="Center",text_align_month="Center",is_bold_year=False,is_bold_month=False,
is_italics_year=False ,is_italics_month=False):
    current_month = start_month
    current_year=current_month.year
    index1=0
    for i in range(no_of_months):
        # Creates Yearly Row
        new_current_year=current_month.year
        if i>0 and new_current_year>current_year and i<no_of_months-1:
            cell = table.cell(0,index1)
            index2=i
            other_cell = table.cell(0,index2-1)
            cell.merge(other_cell)
            index1=index2
            
            cell = table.cell(0,index1)

            if fill_color_year==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_year_r=ImageColor.getcolor(fill_color_year,"RGB")[0]
                    fill_color_year_g=ImageColor.getcolor(fill_color_year,"RGB")[1]
                    fill_color_year_b=ImageColor.getcolor(fill_color_year,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 1 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_year_r,fill_color_year_g,fill_color_year_b)
                except:
                    showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_year} Ex. #FFFF00")
            
            cell.text_frame.clear()

            p=cell.text_frame.paragraphs[0]

            if text_align_year.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_year.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text= str(current_month.year)
            run.font.bold = is_bold_year
            run.font.italic = is_italics_year
            run.font.name = font_name_year
            run.font.size = Pt(font_size_year)
            
            try:
                font_color_year_r=ImageColor.getcolor(font_color_year,"RGB")[0]
                font_color_year_g=ImageColor.getcolor(font_color_year,"RGB")[1]
                font_color_year_b=ImageColor.getcolor(font_color_year,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure Font color for row 1 is selected.")

            try: 
                run.font.color.rgb = RGBColor(font_color_year_r,font_color_year_g,font_color_year_b)
            except:
                showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
                        
            current_year=new_current_year   
        if i==0:
            cell = table.cell(0,i)
            cell.text_frame.clear()

            if fill_color_year==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_year_r=ImageColor.getcolor(fill_color_year,"RGB")[0]
                    fill_color_year_g=ImageColor.getcolor(fill_color_year,"RGB")[1]
                    fill_color_year_b=ImageColor.getcolor(fill_color_year,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 1 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_year_r,fill_color_year_g,fill_color_year_b)
                except:
                    showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_year} Ex. #FFFF00")

            p=cell.text_frame.paragraphs[0]

            if text_align_year.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_year.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT   

            run = p.add_run()
            run.text= str(current_month.year)
            run.font.bold = is_bold_year
            run.font.italic = is_italics_year
            run.font.name = font_name_year
            run.font.size = Pt(font_size_year)
            try:
                font_color_year_r=ImageColor.getcolor(font_color_year,"RGB")[0]
                font_color_year_g=ImageColor.getcolor(font_color_year,"RGB")[1]
                font_color_year_b=ImageColor.getcolor(font_color_year,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure text Font color for row 1 is selected.")

            try: 
                run.font.color.rgb = RGBColor(font_color_year_r,font_color_year_g,font_color_year_b)
            except:
                showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
                       
        if i==no_of_months-1:
            cell = table.cell(0,index1)
            other_cell = table.cell(0,i)
            cell.merge(other_cell)

        # Creates Monthly Table
        cell = table.cell(1,i)
        if fill_color_month==None:
            cell.fill.solid()
            cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            cell.fill.fore_color.brightness = -0.5
        else:
            try:
                fill_color_month_r=ImageColor.getcolor(fill_color_month,"RGB")[0]
                fill_color_month_g=ImageColor.getcolor(fill_color_month,"RGB")[1]
                fill_color_month_b=ImageColor.getcolor(fill_color_month,"RGB")[2]
            except:
                showerror("Text Fill Color Not Selected",
                f"Please ensure text fill color for row 2 is selected.")
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(fill_color_month_r,fill_color_month_g,fill_color_month_b)
            except Exception as e:
                showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_month} Ex. #FFFF00 \n Exception {e}")
        
        cell.text_frame.clear()
        p=cell.text_frame.paragraphs[0]

        if text_align_month.lower()=="right":
            p.alignment = PP_ALIGN.RIGHT
        elif text_align_month.lower()=="center":
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.font.bold = is_bold_month
        run.font.italic = is_italics_month
        run.font.name = font_name_month
        run.text = str(current_month.month)
        # run.text = str(current_month.month_name())[0:3]
        run.font.size = Pt(font_size_month)
        try:
            font_color_month_r=ImageColor.getcolor(font_color_month,"RGB")[0]
            font_color_month_g=ImageColor.getcolor(font_color_month,"RGB")[1]
            font_color_month_b=ImageColor.getcolor(font_color_month,"RGB")[2]
        except:
            showerror("Text Font Color Not Selected",
            f"Please ensure text Font color for row 2 is selected.")
        try: 
            run.font.color.rgb = RGBColor(font_color_month_r,font_color_month_g,font_color_month_b)
        except:
            showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
        
        current_month=current_month+relativedelta(months=1)

#Create a Quarterly Table
def simpleQuarterlyTable(table: pptx.table.Table,start_month: pd._libs.tslibs.timestamps.Timestamp, 
no_of_months: int, font_size_year=12, font_size_month=11, font_color_year="#FFFFFF", font_color_month="#FFFFFF", 
font_name_year="Arial", font_name_month = "Arial", fill_color_year=None, fill_color_month=None,
text_align_year="Center",text_align_month="Center",is_bold_year=False,is_bold_month=False,
is_italics_year=False ,is_italics_month=False):
    current_month = start_month
    current_year = current_month.year
    current_quarter = current_year*10+(current_month.month-1)//3 + 1
    index1=0
    qindex1=0
    for i in range(no_of_months):
        # Creates Yearly Row
        new_current_year=current_month.year
        if i>0 and new_current_year>current_year and i<no_of_months-1:
            cell = table.cell(0,index1)
            index2=i
            other_cell = table.cell(0,index2-1)
            cell.merge(other_cell)
            index1=index2
            
            cell = table.cell(0,index1)

            if fill_color_year==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_year_r=ImageColor.getcolor(fill_color_year,"RGB")[0]
                    fill_color_year_g=ImageColor.getcolor(fill_color_year,"RGB")[1]
                    fill_color_year_b=ImageColor.getcolor(fill_color_year,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 1 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_year_r,fill_color_year_g,fill_color_year_b)
                except:
                    showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_year} Ex. #FFFF00")
            
            cell.text_frame.clear()

            p=cell.text_frame.paragraphs[0]

            if text_align_year.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_year.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text= str(current_month.year)
            run.font.bold = is_bold_year
            run.font.italic = is_italics_year
            run.font.name = font_name_year
            run.font.size = Pt(font_size_year)
            
            try:
                font_color_year_r=ImageColor.getcolor(font_color_year,"RGB")[0]
                font_color_year_g=ImageColor.getcolor(font_color_year,"RGB")[1]
                font_color_year_b=ImageColor.getcolor(font_color_year,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure Font color for row 1 is selected.")

            try: 
                run.font.color.rgb = RGBColor(font_color_year_r,font_color_year_g,font_color_year_b)
            except:
                showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
                        
            current_year=new_current_year   
        if i==0:
            cell = table.cell(0,i)
            cell.text_frame.clear()

            if fill_color_year==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_year_r=ImageColor.getcolor(fill_color_year,"RGB")[0]
                    fill_color_year_g=ImageColor.getcolor(fill_color_year,"RGB")[1]
                    fill_color_year_b=ImageColor.getcolor(fill_color_year,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 1 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_year_r,fill_color_year_g,fill_color_year_b)
                except:
                    showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_year} Ex. #FFFF00")

            p=cell.text_frame.paragraphs[0]

            if text_align_year.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_year.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT   

            run = p.add_run()
            run.text= str(current_month.year)
            run.font.bold = is_bold_year
            run.font.italic = is_italics_year
            run.font.name = font_name_year
            run.font.size = Pt(font_size_year)
            try:
                font_color_year_r=ImageColor.getcolor(font_color_year,"RGB")[0]
                font_color_year_g=ImageColor.getcolor(font_color_year,"RGB")[1]
                font_color_year_b=ImageColor.getcolor(font_color_year,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure Font color for row 1 is selected.")

            try: 
                run.font.color.rgb = RGBColor(font_color_year_r,font_color_year_g,font_color_year_b)
            except:
                showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
                       
        if i==no_of_months-1:
            cell = table.cell(0,index1)
            other_cell = table.cell(0,i)
            cell.merge(other_cell)

        # Create Quarters
        new_current_quarter = new_current_year*10+(current_month.month-1)//3 + 1
        if i>0 and new_current_quarter>current_quarter and i<no_of_months-1:
            cell = table.cell(1,qindex1)
            qindex2=i
            other_cell = table.cell(1,qindex2-1)
            cell.merge(other_cell)
            qindex1=qindex2
            
            cell = table.cell(1,qindex1)

            if fill_color_month==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_month_r=ImageColor.getcolor(fill_color_month,"RGB")[0]
                    fill_color_month_g=ImageColor.getcolor(fill_color_month,"RGB")[1]
                    fill_color_month_b=ImageColor.getcolor(fill_color_month,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 2 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_month_r,fill_color_month_g,fill_color_month_b)
                except Exception as e:
                    showerror("Invalid Fill Color",
                        f"Enter a valid hex color code. Entered Value: {fill_color_month} Ex. #FFFF00 \n Exception {e}")
        
            cell.text_frame.clear()

            p=cell.text_frame.paragraphs[0]

            if text_align_month.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_month.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text= "Q"+str((current_month.month-1)//3+1)
            run.font.bold = is_bold_month
            run.font.italic = is_italics_month
            run.font.name = font_name_month
            run.font.size = Pt(font_size_month)
            
            try:
                font_color_month_r=ImageColor.getcolor(font_color_month,"RGB")[0]
                font_color_month_g=ImageColor.getcolor(font_color_month,"RGB")[1]
                font_color_month_b=ImageColor.getcolor(font_color_month,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure text Font color for row 2 is selected.")
            try: 
                run.font.color.rgb = RGBColor(font_color_month_r,font_color_month_g,font_color_month_b)
            except:
                showerror("Invalid Font Color",
                        "Enter a valid hex color code. Ex. #FFFF00")
                   
            current_quarter=new_current_quarter   
        if i==0:
            cell = table.cell(1,i)
            cell.text_frame.clear()

            if fill_color_month==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_month_r=ImageColor.getcolor(fill_color_month,"RGB")[0]
                    fill_color_month_g=ImageColor.getcolor(fill_color_month,"RGB")[1]
                    fill_color_month_b=ImageColor.getcolor(fill_color_month,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 2 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_month_r,fill_color_month_g,fill_color_month_b)
                except Exception as e:
                    showerror("Invalid Fill Color",
                        f"Enter a valid hex color code. Entered Value: {fill_color_month} Ex. #FFFF00 \n Exception {e}")
        
            p=cell.text_frame.paragraphs[0]

            if text_align_month.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_month.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT   

            run = p.add_run()
            run.text= "Q"+str((current_month.month-1)//3+1)
            run.font.bold = is_bold_month
            run.font.italic = is_italics_month
            run.font.name = font_name_month
            run.font.size = Pt(font_size_month)
            try:
                font_color_month_r=ImageColor.getcolor(font_color_month,"RGB")[0]
                font_color_month_g=ImageColor.getcolor(font_color_month,"RGB")[1]
                font_color_month_b=ImageColor.getcolor(font_color_month,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure text Font color for row 2 is selected.")
            try: 
                run.font.color.rgb = RGBColor(font_color_month_r,font_color_month_g,font_color_month_b)
            except:
                showerror("Invalid Font Color",
                        "Enter a valid hex color code. Ex. #FFFF00")
                   
        if i==no_of_months-1:
            cell = table.cell(1,qindex1)
            other_cell = table.cell(1,i)
            cell.merge(other_cell)
        
        current_month = current_month+relativedelta(months=1)
                        
#Create a Semi Annual Table
def simpleSemiAnnualTable(table: pptx.table.Table,start_month: pd._libs.tslibs.timestamps.Timestamp, 
no_of_months: int, font_size_year=12, font_size_month=11, font_color_year="#FFFFFF", font_color_month="#FFFFFF", 
font_name_year="Arial", font_name_month = "Arial", fill_color_year=None, fill_color_month=None,
text_align_year="Center",text_align_month="Center",is_bold_year=False,is_bold_month=False,
is_italics_year=False ,is_italics_month=False):
    current_month = start_month
    current_year = current_month.year
    current_SemiAnnual = current_year*10+(current_month.month-1)//6 + 1
    index1=0
    qindex1=0
    for i in range(no_of_months):
        # Creates Yearly Row
        new_current_year=current_month.year
        if i>0 and new_current_year>current_year and i<no_of_months-1:
            cell = table.cell(0,index1)
            index2=i
            other_cell = table.cell(0,index2-1)
            cell.merge(other_cell)
            index1=index2
            
            cell = table.cell(0,index1)

            if fill_color_year==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_year_r=ImageColor.getcolor(fill_color_year,"RGB")[0]
                    fill_color_year_g=ImageColor.getcolor(fill_color_year,"RGB")[1]
                    fill_color_year_b=ImageColor.getcolor(fill_color_year,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 1 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_year_r,fill_color_year_g,fill_color_year_b)
                except:
                    showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_year} Ex. #FFFF00")
           
            cell.text_frame.clear()

            p=cell.text_frame.paragraphs[0]

            if text_align_year.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_year.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text= str(current_month.year)
            run.font.bold = is_bold_year
            run.font.italic = is_italics_year
            run.font.name = font_name_year
            run.font.size = Pt(font_size_year)
            
            try:
                font_color_year_r=ImageColor.getcolor(font_color_year,"RGB")[0]
                font_color_year_g=ImageColor.getcolor(font_color_year,"RGB")[1]
                font_color_year_b=ImageColor.getcolor(font_color_year,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure Font color for row 1 is selected.")

            try: 
                run.font.color.rgb = RGBColor(font_color_year_r,font_color_year_g,font_color_year_b)
            except:
                showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
                        
            current_year=new_current_year   
        if i==0:
            cell = table.cell(0,i)
            cell.text_frame.clear()

            if fill_color_year==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_year_r=ImageColor.getcolor(fill_color_year,"RGB")[0]
                    fill_color_year_g=ImageColor.getcolor(fill_color_year,"RGB")[1]
                    fill_color_year_b=ImageColor.getcolor(fill_color_year,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 1 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_year_r,fill_color_year_g,fill_color_year_b)
                except:
                    showerror("Invalid Fill Color",
                    f"Enter a valid hex color code. Entered Value: {fill_color_year} Ex. #FFFF00")
            
            p=cell.text_frame.paragraphs[0]

            if text_align_year.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_year.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT   

            run = p.add_run()
            run.text= str(current_month.year)
            run.font.bold = is_bold_year
            run.font.italic = is_italics_year
            run.font.name = font_name_year
            run.font.size = Pt(font_size_year)
            try:
                font_color_year_r=ImageColor.getcolor(font_color_year,"RGB")[0]
                font_color_year_g=ImageColor.getcolor(font_color_year,"RGB")[1]
                font_color_year_b=ImageColor.getcolor(font_color_year,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure Font color for row 1 is selected.")

            try: 
                run.font.color.rgb = RGBColor(font_color_year_r,font_color_year_g,font_color_year_b)
            except:
                showerror("Invalid Font Color",
                    "Enter a valid hex color code. Ex. #FFFF00")
                       
        if i==no_of_months-1:
            cell = table.cell(0,index1)
            other_cell = table.cell(0,i)
            cell.merge(other_cell)

        # Create SemiAnnuals
        new_current_SemiAnnual = new_current_year*10+(current_month.month-1)//6 + 1
        if i>0 and new_current_SemiAnnual>current_SemiAnnual and i<no_of_months-1:
            cell = table.cell(1,qindex1)
            qindex2=i
            other_cell = table.cell(1,qindex2-1)
            cell.merge(other_cell)
            qindex1=qindex2
            
            cell = table.cell(1,qindex1)

            if fill_color_month==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_month_r=ImageColor.getcolor(fill_color_month,"RGB")[0]
                    fill_color_month_g=ImageColor.getcolor(fill_color_month,"RGB")[1]
                    fill_color_month_b=ImageColor.getcolor(fill_color_month,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 2 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_month_r,fill_color_month_g,fill_color_month_b)
                except Exception as e:
                    showerror("Invalid Fill Color",
                        f"Enter a valid hex color code. Entered Value: {fill_color_month} Ex. #FFFF00 \n Exception {e}")
        
            cell.text_frame.clear()

            p=cell.text_frame.paragraphs[0]

            if text_align_month.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_month.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text= "S"+str((current_month.month-1)//6+1)
            run.font.bold = is_bold_month
            run.font.italic = is_italics_month
            run.font.name = font_name_month
            run.font.size = Pt(font_size_month)
            
            try:
                font_color_month_r=ImageColor.getcolor(font_color_month,"RGB")[0]
                font_color_month_g=ImageColor.getcolor(font_color_month,"RGB")[1]
                font_color_month_b=ImageColor.getcolor(font_color_month,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure text Font color for row 2 is selected.")
            try: 
                run.font.color.rgb = RGBColor(font_color_month_r,font_color_month_g,font_color_month_b)
            except:
                showerror("Invalid Font Color",
                        "Enter a valid hex color code. Ex. #FFFF00")
                        
            current_SemiAnnual=new_current_SemiAnnual   
        if i==0:
            cell = table.cell(1,i)
            cell.text_frame.clear()

            if fill_color_month==None:
                cell.fill.solid()
                cell.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1
            else:
                try:
                    fill_color_month_r=ImageColor.getcolor(fill_color_month,"RGB")[0]
                    fill_color_month_g=ImageColor.getcolor(fill_color_month,"RGB")[1]
                    fill_color_month_b=ImageColor.getcolor(fill_color_month,"RGB")[2]
                except:
                    showerror("Text Fill Color Not Selected",
                    f"Please ensure text fill color for row 2 is selected.")
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(fill_color_month_r,fill_color_month_g,fill_color_month_b)
                except Exception as e:
                    showerror("Invalid Fill Color",
                        f"Enter a valid hex color code. Entered Value: {fill_color_month} Ex. #FFFF00 \n Exception {e}")
        
            p=cell.text_frame.paragraphs[0]

            if text_align_month.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif text_align_month.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT   

            run = p.add_run()
            run.text= "S"+str((current_month.month-1)//6+1)
            run.font.bold = is_bold_month
            run.font.italic = is_italics_month
            run.font.name = font_name_month
            run.font.size = Pt(font_size_month)
            try:
                font_color_month_r=ImageColor.getcolor(font_color_month,"RGB")[0]
                font_color_month_g=ImageColor.getcolor(font_color_month,"RGB")[1]
                font_color_month_b=ImageColor.getcolor(font_color_month,"RGB")[2]
            except:
                showerror("Text Font Color Not Selected",
                f"Please ensure text Font color for row 2 is selected.")
            try: 
                run.font.color.rgb = RGBColor(font_color_month_r,font_color_month_g,font_color_month_b)
            except:
                showerror("Invalid Font Color",
                        "Enter a valid hex color code. Ex. #FFFF00")
                       
        if i==no_of_months-1:
            cell = table.cell(1,qindex1)
            other_cell = table.cell(1,i)
            cell.merge(other_cell)
        
        current_month = current_month+relativedelta(months=1)